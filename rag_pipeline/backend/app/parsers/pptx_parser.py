from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from .base import asset_path_for_storage, detect_languages, make_document, normalize_text
from ..schemas import BlockRecord, CitationAnchor, ImageContext, NormalizedSegment, new_image_id


class PptxParser:
    parser_name = "pptx_parser"

    def __init__(self, extracted_images_dir: Path, root_dir: Path | None = None) -> None:
        self.extracted_images_dir = extracted_images_dir
        self.root_dir = root_dir

    def parse(self, path: Path, relative_path: str):
        deck = Presentation(path)
        document = make_document(path, relative_path, path.suffix.lower().lstrip("."), parser_used=self.parser_name)
        blocks: list[BlockRecord] = []
        segments: list[NormalizedSegment] = []
        images: list[ImageContext] = []
        asset_dir = self.extracted_images_dir / document.document_id
        asset_dir.mkdir(parents=True, exist_ok=True)

        for slide_number, slide in enumerate(deck.slides, start=1):
            title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text else f"Slide {slide_number}"
            slide_anchor = CitationAnchor(slide_number=slide_number, source_label=f"slide {slide_number}")
            slide_section = [document.title, title]
            body_parts: list[str] = []
            image_contexts: list[ImageContext] = []

            blocks.append(
                BlockRecord(
                    block_id=f"{document.document_id}-block-{len(blocks)}",
                    block_type="heading",
                    text=title,
                    order_index=len(blocks),
                    title=title,
                    section_path=slide_section,
                    citation_anchor=slide_anchor,
                    metadata={"slide_number": slide_number, "layout_hint": "title"},
                )
            )

            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = shape.text.strip()
                    if text and text != title:
                        body_parts.append(text)
                image = getattr(shape, "image", None)
                if image is not None:
                    image_name = f"slide-{slide_number}-{new_image_id()}.{image.ext}"
                    image_path = asset_dir / image_name
                    image_path.write_bytes(image.blob)
                    rel_path = asset_path_for_storage(image_path, self.root_dir)
                    image_context = ImageContext(image_id=new_image_id(), image_path=rel_path, related_section_path=title)
                    images.append(image_context)
                    image_contexts.append(image_context)

            notes_text = ""
            try:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                notes_text = ""

            body_text = normalize_text("\n".join(body_parts))
            if body_text or image_contexts:
                blocks.append(
                    BlockRecord(
                        block_id=f"{document.document_id}-block-{len(blocks)}",
                        block_type="paragraph" if body_text else "figure",
                        text=body_text,
                        order_index=len(blocks),
                        title=title,
                        section_path=slide_section,
                        citation_anchor=slide_anchor,
                        image_contexts=image_contexts,
                        metadata={"slide_number": slide_number, "layout_hint": "body"},
                    )
                )
                segments.append(
                    NormalizedSegment(
                        segment_id=f"{document.document_id}-seg-{len(segments)}",
                        segment_type="slide",
                        text=body_text or title,
                        title=title,
                        section_path=slide_section,
                        citation_anchor=slide_anchor,
                        image_contexts=image_contexts,
                        metadata={"slide_number": slide_number, "layout_hint": "body"},
                    )
                )

            notes_clean = normalize_text(notes_text)
            if notes_clean:
                blocks.append(
                    BlockRecord(
                        block_id=f"{document.document_id}-block-{len(blocks)}",
                        block_type="note",
                        text=notes_clean,
                        order_index=len(blocks),
                        title=title,
                        section_path=slide_section,
                        citation_anchor=slide_anchor,
                        metadata={"slide_number": slide_number, "layout_hint": "notes"},
                    )
                )
                segments.append(
                    NormalizedSegment(
                        segment_id=f"{document.document_id}-seg-{len(segments)}",
                        segment_type="slide",
                        text=notes_clean,
                        title=title,
                        section_path=slide_section,
                        citation_anchor=slide_anchor,
                        metadata={"slide_number": slide_number, "layout_hint": "notes"},
                    )
                )

        document.images = images
        document.blocks = blocks
        document.segments = segments
        document.metadata = {
            **document.metadata,
            "document_structure": "slides",
            "chunking_hints": _build_chunking_hints(document),
        }
        document.language_tags = detect_languages(block.text for block in blocks if block.text)
        for block in document.blocks:
            block.language_tags = document.language_tags
        for segment in document.segments:
            segment.language_tags = document.language_tags
        return document


def _build_chunking_hints(document) -> dict[str, object]:
    slides: list[dict[str, object]] = []
    blocks_by_slide: dict[int, list[BlockRecord]] = {}
    for block in document.blocks:
        slide_number = block.citation_anchor.slide_number
        if slide_number is None:
            continue
        blocks_by_slide.setdefault(slide_number, []).append(block)

    for slide_number in sorted(blocks_by_slide):
        slide_blocks = blocks_by_slide[slide_number]
        section_path = slide_blocks[0].section_path if slide_blocks else [document.title, f"Slide {slide_number}"]
        layout_hints: list[str] = []
        block_ids: list[str] = []
        for block in slide_blocks:
            block_ids.append(block.block_id)
            layout_hint = str(block.metadata.get("layout_hint", "")).strip()
            if layout_hint and layout_hint not in layout_hints:
                layout_hints.append(layout_hint)
        slides.append(
            {
                "slide_number": slide_number,
                "section_path": section_path,
                "block_ids": block_ids,
                "layout_hints": layout_hints,
            }
        )

    return {
        "preferred_strategy": "slide_window",
        "structure": "slides",
        "slide_count": len(slides),
        "slides": slides,
    }
