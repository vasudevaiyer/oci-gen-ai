from pathlib import Path

from pptx import Presentation

from backend.app.parsers.pptx_parser import PptxParser


def test_pptx_parser_emits_slide_structure_hints(tmp_path: Path) -> None:
    pptx_path = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Remote Work"
    slide.placeholders[1].text = "Employees may work from home up to 30 days per year."
    slide.notes_slide.notes_text_frame.text = "Manager approval is required."
    presentation.save(pptx_path)

    document = PptxParser(tmp_path / "images", tmp_path).parse(pptx_path, "deck.pptx")

    assert document.metadata["document_structure"] == "slides"
    assert document.metadata["chunking_hints"]["preferred_strategy"] == "slide_window"
    assert document.metadata["chunking_hints"]["slide_count"] == 1
    assert any(block.block_type == "heading" for block in document.blocks)
    assert any(block.metadata.get("layout_hint") == "body" for block in document.blocks)
    assert any(block.metadata.get("layout_hint") == "notes" for block in document.blocks)
